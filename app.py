import streamlit as st
import pandas as pd
import json
import io
import re
import anthropic
from openpyxl import Workbook
from openpyxl.styles import PatternFill, Font, Alignment
from openpyxl.utils import get_column_letter
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

st.set_page_config(page_title="Doctor Response Analyzer", page_icon="⚕️", layout="wide")

COLORS = ["#1a2e5a","#e5333a","#f7a826","#16a34a","#7c3aed","#0891b2","#be185d","#b45309","#065f46","#1d4ed8","#9d174d","#0f766e"]
PALH   = [c.replace("#","") for c in COLORS]

st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Syne:wght@700;800&family=Inter:wght@400;500;600&display=swap');
html,body,[class*="css"]{font-family:'Inter',sans-serif}
.banner{background:#1a2e5a;padding:1.6rem 2rem;border-radius:12px;margin-bottom:1.5rem;color:white}
.banner h1{font-family:'Syne',sans-serif;font-size:1.7rem;font-weight:800;margin:0 0 4px}
.banner p{margin:0;opacity:.6;font-size:13px}
.rtag{display:inline-block;background:#e5333a;color:white;font-size:11px;font-weight:600;padding:2px 10px;border-radius:20px;margin-bottom:8px}
.qb{border-left:3px solid #e5333a;border-radius:0 8px 8px 0;padding:9px 13px;margin:7px 0;background:#f5f4f1;font-size:13px;font-style:italic;color:#555;line-height:1.55}
.qbk{font-style:normal;font-size:10px;font-weight:700;color:#e5333a;margin-top:3px;text-transform:uppercase;letter-spacing:.4px}
</style>
""", unsafe_allow_html=True)

st.markdown("""
<div class="banner">
  <div class="rtag">⚕ AI Research Tool</div>
  <h1>Doctor Response Analyzer</h1>
  <p>Upload responses · AI bucketing · Excel + PowerPoint export · One link for your whole team</p>
</div>
""", unsafe_allow_html=True)

# API key
if "ANTHROPIC_API_KEY" in st.secrets:
    api_key = st.secrets["ANTHROPIC_API_KEY"]
else:
    api_key = st.sidebar.text_input("Anthropic API key", type="password", placeholder="sk-ant-...")
    if not api_key:
        st.sidebar.warning("Add API key to run analysis")

defaults = {"buckets":[],"total":0,"tagged_df":None,"col":None,"bg":"","question":"","done":False,"mode":""}
for k,v in defaults.items():
    if k not in st.session_state: st.session_state[k] = v


def extract_doctor_lines(text):
    """
    Smart extraction that handles 3 formats:
    1. Full transcript: SPEAKER_A/SPEAKER_B turns -> extracts only SPEAKER_B
    2. Single doctor response (plain text) -> returns as-is
    3. AI Moderator / Doctor format -> extracts Doctor lines
    """
    text = str(text).strip()
    if not text or len(text) < 10:
        return []

    # Detect if this is a transcript (contains speaker tags)
    is_transcript = bool(re.search(
        r'SPEAKER_[AB]:|AI Moderator:|Doctor:|Moderator:|\[GT\]', text, re.I))

    if not is_transcript:
        # Plain response — use as-is if long enough
        if len(text) > 15:
            return [text]
        return []

    # Split on any speaker label
    parts = re.split(r'(SPEAKER_[AB]\s*:|AI Moderator\s*:|Doctor\s*:|Moderator\s*:|Interviewer\s*:)', text)
    current_speaker = None
    results = []

    for part in parts:
        part = part.strip()
        if not part:
            continue
        if re.match(r'SPEAKER_A\s*:|AI Moderator\s*:|Moderator\s*:|Interviewer\s*:', part, re.I):
            current_speaker = 'A'
        elif re.match(r'SPEAKER_B\s*:|Doctor\s*:', part, re.I):
            current_speaker = 'B'
        else:
            if current_speaker == 'B':
                clean = re.sub(r'\[GT\]', '', part).strip()
                clean = re.sub(r'\s+', ' ', clean).strip()
                # Filter out very short / non-substantive replies
                skip_patterns = [
                    r'^(yes|no|yeah|nope|ok|okay|sure|right|correct|that\'?s? (it|correct|right)|nothing else|not really|n/a)\.*$',
                    r'^(are you still there|hello|thank you)\?*$',
                    r'^\d+$',
                ]
                if len(clean) > 20 and not any(re.match(p, clean, re.I) for p in skip_patterns):
                    results.append(clean)
    return results


def process_column(df, col):
    """
    Process a DataFrame column.
    Returns (responses, row_map) where row_map[response_idx] = original_row_idx
    Handles both single-response rows and transcript rows.
    """
    responses = []
    row_map = []   # response index -> original row index

    for row_idx, val in enumerate(df[col].astype(str)):
        lines = extract_doctor_lines(val)
        for line in lines:
            responses.append(line)
            row_map.append(row_idx)

    return responses, row_map


# ── Step 1 ────────────────────────────────────────────────────────────────────
with st.expander("**Step 1 — Study context**", expanded=not st.session_state.done):
    c1,c2 = st.columns(2)
    bg       = c1.text_area("Study background", height=90,
                 placeholder="e.g. Post-consultation survey with oncologists discussing VORANIGO prescribing barriers...")
    question = c2.text_input("Question asked to doctors",
                 placeholder="e.g. What are the primary reasons you choose Voranigo for IDH mutated Glioma patients?")

# ── Step 2 ────────────────────────────────────────────────────────────────────
st.markdown("### 📂 Step 2 — Upload Excel or CSV")
uploaded = st.file_uploader("Upload file", type=["xlsx","xls","csv"], label_visibility="collapsed")

df, sel_col = None, None
if uploaded:
    try:
        df = pd.read_csv(uploaded, dtype=str).fillna("") if uploaded.name.endswith(".csv") \
             else pd.read_excel(uploaded, dtype=str).fillna("")
        st.success(f"✓ **{uploaded.name}** — {len(df)} rows, {len(df.columns)} columns")
        kw = ["response","answer","comment","text","doctor","feedback","reply","verbatim","transcript","column"]
        auto = next((c for c in df.columns if any(k in c.lower() for k in kw)), df.columns[0])
        sel_col = st.selectbox("Which column has the responses/transcripts?",
                               df.columns.tolist(), index=df.columns.tolist().index(auto))
        if sel_col and len(df) > 0:
            # Preview with extraction info
            preview_responses, _ = process_column(df, sel_col)
            st.info(f"📋 Detected **{preview_responses and len(preview_responses) or 0} doctor statements** "
                    f"across {len(df)} rows "
                    f"({'transcripts detected — will extract SPEAKER_B lines automatically' if any(re.search(r'SPEAKER_[AB]:', str(v)) for v in df[sel_col].head(3)) else 'single responses per row'})")
            if preview_responses:
                st.caption(f'↳ First extracted statement: *"{preview_responses[0][:180]}…"*')
    except Exception as e:
        st.error(f"Cannot read file: {e}")

# ── Step 3 ────────────────────────────────────────────────────────────────────
st.markdown("### 🤖 Step 3 — Run AI Analysis")

if st.button("▶ Run AI Analysis",
             disabled=not (df is not None and sel_col and api_key),
             type="primary"):

    with st.spinner("Extracting doctor statements…"):
        responses, row_map = process_column(df, sel_col)

    if len(responses) < 3:
        st.error(f"Only {len(responses)} doctor statements found. "
                 f"Check your column selection — if this is a transcript file, "
                 f"make sure SPEAKER_B or Doctor: labels are present.")
    else:
        st.info(f"✓ Extracted **{len(responses)} doctor statements** — sending to AI for bucketing")

        with st.spinner(f"Analyzing {len(responses)} statements… 20–40 seconds"):
            sys_p = """You are an expert qualitative researcher for medical/pharma studies.
Analyse the doctor statements and return ONLY valid JSON — no markdown, no backticks, no explanation.

{
  "totalResponses": <integer — total number of statements analyzed>,
  "buckets": [
    {
      "name": "<4-7 word insight-forward label>",
      "count": <integer>,
      "percentage": <float 1dp>,
      "theme": "<one sentence core insight>",
      "quotes": ["<verbatim 10-25 word fragment from actual doctor text>","<another>","<third>"],
      "responseIndices": [<0-based integers — index within the statements list>]
    }
  ]
}

CRITICAL RULES:
- Create exactly 10-12 buckets
- Every index 0..(totalResponses-1) must appear in EXACTLY ONE bucket — zero uncategorized
- counts must sum to totalResponses
- percentages must sum to 100.0
- Order by count descending
- Names must be specific insight-forward labels (e.g. "PFS as Primary Efficacy Endpoint" not just "Efficacy")
- Quotes must be verbatim text from the actual statements provided"""

            # Send up to 200 statements to AI (most representative sample)
            sample = responses[:200]
            user_msg = (
                f"Background: {bg or 'Physician research on IDH mutated Glioma treatment'}\n"
                f"Question: {question or 'Primary reasons for choosing Voranigo'}\n\n"
                f"Doctor statements ({len(sample)} total, indices 0–{len(sample)-1}):\n"
                + "\n".join(f"[{i}] {r}" for i,r in enumerate(sample))
            )

            try:
                client = anthropic.Anthropic(api_key=api_key)
                msg = client.messages.create(
                    model="claude-sonnet-4-5",
                    max_tokens=2000,
                    system=sys_p,
                    messages=[{"role":"user","content":user_msg}]
                )
                raw = re.sub(r'^```json\s*','',msg.content[0].text.strip(),flags=re.I)
                raw = re.sub(r'^```','',raw).replace('```','').strip()
                parsed = json.loads(raw)

                buckets   = parsed["buckets"]
                total_r   = parsed["totalResponses"]

                # Build tagged dataframe — one row per original Excel row
                # A row gets the bucket of its most representative statement
                tagged = df.copy()
                tagged["Bucket"] = "N/A"
                tagged["Doctor Statements Extracted"] = ""

                # Map filtered response index -> bucket name
                resp_to_bucket = {}
                for b in buckets:
                    for fi in (b.get("responseIndices") or []):
                        if fi < len(responses):
                            resp_to_bucket[fi] = b["name"]

                # Assign bucket to each original row
                # If a row has multiple statements, use the first bucketed one
                row_statements = {}   # row_idx -> [statements]
                row_bucket     = {}   # row_idx -> bucket name
                for resp_idx, orig_row in enumerate(row_map):
                    if orig_row not in row_statements:
                        row_statements[orig_row] = []
                    row_statements[orig_row].append(responses[resp_idx])
                    if orig_row not in row_bucket and resp_idx in resp_to_bucket:
                        row_bucket[orig_row] = resp_to_bucket[resp_idx]

                for row_idx in range(len(tagged)):
                    tagged.at[row_idx, "Bucket"] = row_bucket.get(row_idx, "N/A")
                    stmts = row_statements.get(row_idx, [])
                    tagged.at[row_idx, "Doctor Statements Extracted"] = " | ".join(stmts[:3])

                st.session_state.update({
                    "buckets": buckets,
                    "total": total_r,
                    "tagged_df": tagged,
                    "col": sel_col,
                    "bg": bg,
                    "question": question,
                    "done": True,
                    "mode": "transcript" if len(responses) > len(df) else "single"
                })
                st.rerun()

            except Exception as e:
                st.error(f"Analysis failed: {e}")
                st.exception(e)

# ── Results ───────────────────────────────────────────────────────────────────
if st.session_state.done and st.session_state.buckets:
    B         = st.session_state.buckets
    total     = st.session_state.total
    tagged_df = st.session_state.tagged_df

    st.divider()
    st.markdown("## Results")
    m1,m2,m3,m4 = st.columns(4)
    m1.metric("Doctor Statements", total)
    m2.metric("Buckets", len(B))
    m3.metric("Top Bucket", f"{B[0]['percentage']}%")
    m4.metric("Major Themes ≥10%", sum(1 for b in B if b["percentage"]>=10))

    t1,t2,t3,t4 = st.tabs(["📋 Table","📊 Visualize","💬 Quotes","🏷️ Tagged"])

    with t1:
        st.dataframe(pd.DataFrame([{
            "#": i+1, "Bucket": b["name"], "Count": b["count"],
            "%": f"{b['percentage']}%", "Theme": b["theme"],
            "Sample quote": b["quotes"][0] if b.get("quotes") else ""
        } for i,b in enumerate(B)]), use_container_width=True, hide_index=True)

    with t2:
        viz = st.selectbox("Choose style", [
            "Horizontal bar chart + quotes","Stat cards per bucket",
            "Ranked driver list","Donut chart","Summary table"
        ], key="viz_pick")
        names  = [b["name"] for b in B]
        counts = [b["count"] for b in B]
        pcts   = [b["percentage"] for b in B]
        colors = [COLORS[i%len(COLORS)] for i in range(len(B))]

        if "bar" in viz:
            fig = go.Figure(go.Bar(x=counts, y=names, orientation='h', marker_color=colors,
                text=[f"{p}% ({c})" for c,p in zip(counts,pcts)], textposition='outside'))
            fig.update_layout(height=max(300,len(B)*42+100), margin=dict(l=20,r=140,t=10,b=20),
                yaxis=dict(autorange="reversed"), plot_bgcolor='white', paper_bgcolor='white',
                xaxis_title="Number of statements", font=dict(family="Inter",size=12))
            fig.update_xaxes(showgrid=True, gridcolor='#f0f0f0')
            st.plotly_chart(fig, use_container_width=True)
            st.markdown("**Key quotes:**")
            for b in B[:3]:
                if b.get("quotes"):
                    st.markdown(f'<div class="qb">"{b["quotes"][0]}"<div class="qbk">{b["name"]} · {b["percentage"]}%</div></div>',unsafe_allow_html=True)

        elif "cards" in viz:
            cols = st.columns(3)
            for i,b in enumerate(B):
                c = COLORS[i%len(COLORS)]
                q = f'<div style="font-size:11px;font-style:italic;color:#888;border-left:2px solid {c};padding-left:7px;margin-top:6px">"{b["quotes"][0]}"</div>' if b.get("quotes") else ""
                with cols[i%3]:
                    st.markdown(f'<div style="background:white;border:1px solid rgba(0,0,0,0.08);border-radius:10px;padding:13px;border-left:4px solid {c};margin-bottom:8px"><div style="font-size:26px;font-weight:800;color:{c};line-height:1;margin-bottom:4px">{b["percentage"]}%</div><div style="font-size:13px;font-weight:600;margin-bottom:3px">{b["name"]}</div><div style="font-size:12px;color:#6b6b78;margin-bottom:6px">{b["theme"]}</div>{q}</div>',unsafe_allow_html=True)

        elif "ranked" in viz:
            for i,b in enumerate(B):
                c = COLORS[i%len(COLORS)]
                impact = "🔴 Most impactful" if i<int(len(B)*0.3) else "🟡 Moderate" if i<int(len(B)*0.7) else "🟢 Least impactful"
                if i==0 or i==int(len(B)*0.3) or i==int(len(B)*0.7):
                    st.markdown(f"**{impact}**")
                q = f'<div style="font-size:11px;font-style:italic;color:#888;margin-top:5px;border-left:2px solid #e0e0e8;padding-left:7px">"{b["quotes"][0]}"</div>' if b.get("quotes") else ""
                st.markdown(f'<div style="display:flex;align-items:flex-start;gap:12px;padding:11px 14px;background:white;border:1px solid rgba(0,0,0,0.08);border-radius:10px;margin-bottom:6px;border-left:4px solid {c}"><div style="font-size:18px;font-weight:800;color:{c};min-width:26px;text-align:center">{i+1}</div><div style="flex:1"><div style="font-size:13px;font-weight:600">{b["name"]}</div><div style="font-size:11px;color:#6b6b78">{b["theme"]}</div>{q}</div><div style="font-size:20px;font-weight:800;color:{c}">{b["percentage"]}%</div></div>',unsafe_allow_html=True)

        elif "donut" in viz:
            fig = go.Figure(go.Pie(labels=names, values=counts, hole=0.52,
                marker_colors=colors, textinfo='percent'))
            fig.update_layout(height=420, margin=dict(l=20,r=20,t=10,b=20),
                showlegend=True, paper_bgcolor='white')
            st.plotly_chart(fig, use_container_width=True)

        else:
            for i,b in enumerate(B):
                c = COLORS[i%len(COLORS)]
                q = f'<em>"{b["quotes"][0]}"</em>' if b.get("quotes") else "—"
                st.markdown(f'<div style="display:grid;grid-template-columns:auto 1fr auto;gap:12px;align-items:center;padding:10px 12px;background:white;border:1px solid rgba(0,0,0,0.07);border-radius:8px;margin-bottom:5px;border-left:4px solid {c}"><div style="font-size:14px;font-weight:800;color:{c};min-width:28px;text-align:center">{i+1}</div><div><div style="font-size:13px;font-weight:600">{b["name"]}</div><div style="font-size:11px;color:#888;margin-top:2px">{b["theme"]}</div><div style="font-size:11px;color:#aaa;margin-top:3px">{q}</div></div><div style="font-size:22px;font-weight:800;color:{c}">{b["percentage"]}%</div></div>',unsafe_allow_html=True)

    with t3:
        for b in B:
            for q in (b.get("quotes") or []):
                st.markdown(f'<div class="qb">"{q}"<div class="qbk">{b["name"]} · {b["percentage"]}%</div></div>',unsafe_allow_html=True)

    with t4:
        na_count = (tagged_df["Bucket"]=="N/A").sum()
        if na_count > 0:
            st.info(f"{na_count} rows marked N/A — these rows had no extractable doctor statements (e.g. metadata, empty rows).")
        st.dataframe(tagged_df, use_container_width=True, height=420)

    # ── Export ────────────────────────────────────────────────────────────────
    st.divider()
    st.markdown("## 📥 Export")
    e1, e2 = st.columns(2)

    with e1:
        st.markdown("#### 📊 Excel")
        if st.button("Generate Excel", type="primary", use_container_width=True):
            with st.spinner("Building Excel…"):
                wb = Workbook()
                ws1 = wb.active; ws1.title = "Responses (Bucketed)"
                for j,h in enumerate(tagged_df.columns,1):
                    c = ws1.cell(row=1,column=j,value=h)
                    c.font = Font(bold=True,color="FFFFFF",name="Calibri",size=10)
                    c.fill = PatternFill("solid",fgColor="1A2E5A")
                    c.alignment = Alignment(horizontal="center")
                for i,row in tagged_df.iterrows():
                    for j,val in enumerate(row,1):
                        cell = ws1.cell(row=i+2,column=j,value=str(val))
                        cell.font = Font(name="Calibri",size=10)
                        cell.alignment = Alignment(wrap_text=True,vertical="top")
                        if tagged_df.columns[j-1] == "Bucket":
                            bkt = next((b for b in B if b["name"]==str(val)),None)
                            if bkt:
                                hx = PALH[B.index(bkt)%len(PALH)]
                                cell.fill = PatternFill("solid",fgColor=hx+"22" if len(hx)==6 else "EEEEEE")
                                cell.font = Font(name="Calibri",size=10,color=hx,bold=True)
                for cc in ws1.columns:
                    ws1.column_dimensions[get_column_letter(cc[0].column)].width = min(
                        max(len(str(c.value or "")) for c in cc)+2, 60)

                ws2 = wb.create_sheet("Bucket Summary")
                for j,h in enumerate(["Rank","Bucket","Count","%","Theme","Quote 1","Quote 2","Quote 3"],1):
                    c = ws2.cell(row=1,column=j,value=h)
                    c.font = Font(bold=True,color="FFFFFF",name="Calibri",size=10)
                    c.fill = PatternFill("solid",fgColor="E5333A")
                    c.alignment = Alignment(horizontal="center")
                for i,b in enumerate(B):
                    qs = b.get("quotes",[])
                    for j,v in enumerate([i+1,b["name"],b["count"],f"{b['percentage']}%",b["theme"],
                        qs[0] if len(qs)>0 else "",qs[1] if len(qs)>1 else "",qs[2] if len(qs)>2 else ""],1):
                        c = ws2.cell(row=i+2,column=j,value=v)
                        c.font = Font(name="Calibri",size=10)
                        c.alignment = Alignment(wrap_text=True,vertical="top")
                for col_letter,width in zip(["A","B","C","D","E","F","G","H"],[6,38,8,10,55,50,50,50]):
                    ws2.column_dimensions[col_letter].width = width

                buf = io.BytesIO(); wb.save(buf); buf.seek(0)
                st.download_button("⬇️ Download Excel", data=buf.getvalue(),
                    file_name="doctor_responses_bucketed.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    use_container_width=True)

    with e2:
        st.markdown("#### 📑 PowerPoint — 4 slides")
        st.caption("Slide 1: Bold infographic · Slide 2: Numbered steps · Slide 3: Bar chart + quotes · Slide 4: Full table")
        if st.button("Generate PowerPoint", type="primary", use_container_width=True):
            with st.spinner("Building 4-slide deck…"):

                prs = Presentation()
                prs.slide_width  = Inches(13.33)
                prs.slide_height = Inches(7.5)
                blank = prs.slide_layouts[6]

                def bgs(slide, h):
                    slide.background.fill.solid()
                    slide.background.fill.fore_color.rgb = RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))

                def lighten(h, f=0.85):
                    h=h[:6]; r,g,b=int(h[0:2],16),int(h[2:4],16),int(h[4:6],16)
                    return f'{min(int(r+(255-r)*f),255):02X}{min(int(g+(255-g)*f),255):02X}{min(int(b+(255-b)*f),255):02X}'

                def R(slide,x,y,w,h,c):
                    c=str(c)[:6]
                    sp=slide.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(h))
                    sp.fill.solid(); sp.fill.fore_color.rgb=RGBColor.from_string(c)
                    sp.line.fill.background(); return sp

                def oval(slide,x,y,w,h,c):
                    c=str(c)[:6]
                    sp=slide.shapes.add_shape(9,Inches(x),Inches(y),Inches(w),Inches(h))
                    sp.fill.solid(); sp.fill.fore_color.rgb=RGBColor.from_string(c)
                    sp.line.fill.background(); return sp

                def T(slide,text,x,y,w,h,sz=12,bold=False,col="000000",italic=False,align=PP_ALIGN.LEFT):
                    col=str(col)[:6]
                    tb=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
                    tb.word_wrap=True; tf=tb.text_frame; tf.word_wrap=True
                    p=tf.paragraphs[0]; p.alignment=align
                    run=p.add_run(); run.text=str(text)
                    run.font.size=Pt(sz); run.font.bold=bold; run.font.italic=italic
                    run.font.color.rgb=RGBColor.from_string(col)

                def add_connector(slide,x1,y1,x2,y2,c,width_pt=1.0):
                    cn=slide.shapes.add_connector(1,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
                    cn.line.color.rgb=RGBColor.from_string(str(c)[:6])
                    cn.line.width=Pt(width_pt)

                q_txt = st.session_state.question or "Key Themes from Doctor Responses"
                bg_txt = st.session_state.bg or ""

                # ── SLIDE 1: Bold stat infographic ───────────────────────────
                s1=prs.slides.add_slide(blank); bgs(s1,"0F0F1E")
                R(s1,0,0,0.18,7.5,"E5333A")
                T(s1,"PHYSICIAN RESEARCH INSIGHTS",0.38,0.28,9,0.32,sz=9,bold=True,col="E5333A")
                T(s1,q_txt,0.38,0.65,9.5,0.95,sz=26,bold=True,col="FFFFFF")
                if bg_txt:
                    T(s1,bg_txt[:110],0.38,1.7,9,0.5,sz=11,italic=True,col="8888AA")

                card_cols = ["E5333A","1A2E5A","F7A826","16A34A","7C3AED"]
                top5 = B[:5]; cw,ch = 2.42,2.1; gap=0.14
                total_w = len(top5)*(cw+gap)-gap
                sx = (13.33-total_w)/2
                for i,b in enumerate(top5):
                    bx=sx+i*(cw+gap); by=2.38; ch_hex=card_cols[i%len(card_cols)]
                    R(s1,bx,by,cw,ch,ch_hex)
                    T(s1,f"{b['percentage']}%",bx,by+0.1,cw,0.75,sz=36,bold=True,col="FFFFFF",align=PP_ALIGN.CENTER)
                    R(s1,bx+cw/2-0.42,by+0.85,0.84,0.26,"FFFFFF")
                    T(s1,f"n = {b['count']}",bx+cw/2-0.42,by+0.85,0.84,0.26,sz=9,bold=True,col=ch_hex,align=PP_ALIGN.CENTER)
                    nm=b["name"][:30]+("…" if len(b["name"])>30 else "")
                    T(s1,nm,bx+0.1,by+1.18,cw-0.2,0.68,sz=10,bold=True,col="FFFFFF",align=PP_ALIGN.CENTER)

                rest=B[5:12]
                if rest:
                    rw=(13.33-0.56)/len(rest)
                    for i,b in enumerate(rest):
                        rx=0.38+i*rw; ch_hex=PALH[(i+5)%len(PALH)]
                        R(s1,rx,4.72,rw-0.1,1.62,ch_hex)
                        T(s1,f"{b['percentage']}%",rx,4.77,rw-0.1,0.55,sz=20,bold=True,col="FFFFFF",align=PP_ALIGN.CENTER)
                        nm2=b["name"][:20]+("…" if len(b["name"])>20 else "")
                        T(s1,nm2,rx+0.04,5.32,rw-0.18,0.78,sz=8,col="CCCCEE",align=PP_ALIGN.CENTER)

                T(s1,f"n = {total} physician statements  ·  {len(B)} themes identified",
                  0,6.9,13.33,0.3,sz=9,col="444460",align=PP_ALIGN.CENTER)

                # ── SLIDE 2: Numbered steps infographic ──────────────────────
                s2=prs.slides.add_slide(blank); bgs(s2,"FFFFFF")
                T(s2,"TOP THEMES — RANKED",0.5,0.2,12,0.3,sz=8,bold=True,col="E5333A")
                T(s2,q_txt,0.5,0.52,8.5,0.48,sz=18,bold=True,col="0F0F18")

                step_colors=["E5333A","F7A826","1A2E5A","16A34A","7C3AED"]
                cr=0.42; cx_circ=0.65; cx_box=1.36; bw=5.6; bh=0.95; sy=1.18; sg=1.12

                for i,b in enumerate(B[:5]):
                    ch_hex=step_colors[i%len(step_colors)]; cy=sy+i*sg
                    if i<4:
                        add_connector(s2,cx_circ+cr/2,cy+cr,cx_circ+cr/2,cy+sg,"DDDDDD",0.8)
                    oval(s2,cx_circ,cy,cr,cr,ch_hex)
                    T(s2,str(i+1),cx_circ,cy,cr,cr,sz=16,bold=True,col="FFFFFF",align=PP_ALIGN.CENTER)
                    R(s2,cx_box,cy-0.04,bw,bh,lighten(ch_hex,0.88))
                    R(s2,cx_box,cy-0.04,0.06,bh,ch_hex)
                    T(s2,b["name"],cx_box+0.16,cy,bw-0.45,0.35,sz=11,bold=True,col="0F0F18")
                    T(s2,b["theme"][:62],cx_box+0.16,cy+0.34,bw-0.45,0.34,sz=9,italic=True,col="555566")
                    R(s2,cx_box+bw+0.12,cy+0.12,0.82,0.52,ch_hex)
                    T(s2,f"{b['percentage']}%",cx_box+bw+0.12,cy+0.12,0.82,0.52,sz=13,bold=True,col="FFFFFF",align=PP_ALIGN.CENTER)

                # Right panel — remaining buckets
                rx_p=8.16
                T(s2,"OTHER THEMES",rx_p,1.1,4.8,0.3,sz=8,bold=True,col="E5333A")
                rest_s2=B[5:]; mh=min(0.6,5.5/max(len(rest_s2),1))
                for i,b in enumerate(rest_s2):
                    ch_hex=PALH[(i+5)%len(PALH)]; my=1.5+i*mh
                    R(s2,rx_p,my,4.78,mh-0.06,lighten(ch_hex,0.92))
                    R(s2,rx_p,my,0.06,mh-0.06,ch_hex)
                    T(s2,b["name"][:36],rx_p+0.14,my+0.03,3.4,mh-0.1,sz=9,bold=True,col="222233")
                    T(s2,f"{b['percentage']}%",rx_p+3.56,my+0.03,1.0,mh-0.1,sz=11,bold=True,col=ch_hex,align=PP_ALIGN.RIGHT)

                # ── SLIDE 3: Bar chart + quotes ──────────────────────────────
                s3=prs.slides.add_slide(blank); bgs(s3,"FFFFFF")
                T(s3,"RESPONSE DISTRIBUTION",0.5,0.22,8,0.28,sz=8,bold=True,col="E5333A")
                T(s3,q_txt,0.5,0.52,12.5,0.48,sz=17,bold=True,col="0F0F18")
                T(s3,f"All {len(B)} themes  ·  n = {total} statements",0.5,1.02,8,0.26,sz=9,italic=True,col="888899")

                n=len(B); rh=5.3/n; max_c=max(b["count"] for b in B) or 1
                for i,b in enumerate(B):
                    ch_hex=PALH[i%len(PALH)]; bw2=(b["count"]/max_c)*5.6; yp=1.38+i*rh
                    T(s3,b["name"][:36],0.5,yp+0.01,3.4,rh-0.04,sz=8,col="222233")
                    R(s3,4.0,yp+rh*0.18,5.8,rh*0.6,"F4F3F0")
                    if bw2>0.05: R(s3,4.0,yp+rh*0.18,bw2,rh*0.6,ch_hex)
                    T(s3,f"{b['percentage']}%  n={b['count']}",4.08+bw2,yp+0.01,2.2,rh-0.04,sz=8,bold=True,col=ch_hex)

                R(s3,10.08,1.38,3.08,5.82,"F7F6F3")
                T(s3,"KEY QUOTES",10.22,1.5,2.8,0.24,sz=7,bold=True,col="E5333A")
                top_qs=[(b,b["quotes"][0]) for b in B[:3] if b.get("quotes")]
                for idx,(b,q) in enumerate(top_qs[:3]):
                    ch_hex=PALH[idx%len(PALH)]; qy=1.88+idx*1.38
                    R(s3,10.22,qy,0.05,0.85,ch_hex)
                    T(s3,f'"{q[:88]}{"…" if len(q)>88 else ""}"',10.32,qy+0.02,2.6,0.68,sz=8,italic=True,col="333344")
                    T(s3,f"— {b['name'][:26]}",10.32,qy+0.7,2.6,0.2,sz=7,bold=True,col=ch_hex)

                # ── SLIDE 4: Full table ───────────────────────────────────────
                s4=prs.slides.add_slide(blank); bgs(s4,"F7F6F3")
                T(s4,"FULL BREAKDOWN",0.5,0.22,9,0.28,sz=8,bold=True,col="E5333A")
                T(s4,f"All {len(B)} themes ranked by frequency",0.5,0.54,9,0.38,sz=15,bold=True,col="0F0F18")
                T(s4,f"n = {total} physician statements",0.5,0.96,5,0.26,sz=9,italic=True,col="888899")

                cxs=[0.5,0.98,3.88,4.56,5.18]; cws=[0.44,2.86,0.64,0.58,7.54]
                hy=1.3
                for lbl,cx,cw2 in zip(["#","Bucket","Count","%","Core theme"],cxs,cws):
                    R(s4,cx,hy,cw2,0.36,"1A2E5A")
                    T(s4,lbl,cx+0.05,hy+0.04,cw2-0.1,0.28,sz=9,bold=True,col="FFFFFF")

                trh=min(0.37,(7.5-hy-0.5)/(len(B)+1))
                for i,b in enumerate(B):
                    ry=hy+0.36+i*trh; bg_c="FFFFFF" if i%2==0 else "F2F1EE"; ch_hex=PALH[i%len(PALH)]
                    for cx,cw2 in zip(cxs,cws): R(s4,cx,ry,cw2,trh,bg_c)
                    T(s4,str(i+1),cxs[0]+0.05,ry+0.03,0.34,trh-0.06,sz=9,col="888888")
                    T(s4,b["name"],cxs[1]+0.05,ry+0.03,cws[1]-0.1,trh-0.06,sz=9,bold=True,col=ch_hex)
                    T(s4,str(b["count"]),cxs[2]+0.05,ry+0.03,cws[2]-0.1,trh-0.06,sz=9,col="333344")
                    T(s4,f"{b['percentage']}%",cxs[3]+0.05,ry+0.03,cws[3]-0.1,trh-0.06,sz=9,bold=True,col=ch_hex)
                    T(s4,b["theme"][:85],cxs[4]+0.05,ry+0.03,cws[4]-0.1,trh-0.06,sz=8,italic=True,col="555566")

                buf2=io.BytesIO(); prs.save(buf2); buf2.seek(0)
                st.download_button("⬇️ Download PowerPoint", data=buf2.getvalue(),
                    file_name="doctor_response_analysis.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True)

    st.divider()
    if st.button("↺ New analysis"):
        for k,v in defaults.items(): st.session_state[k] = v
        st.rerun()

st.markdown("---")
st.markdown("<div style='text-align:center;font-size:11px;color:#aaa'>Doctor Response Analyzer · Powered by Claude · No data stored</div>",unsafe_allow_html=True)
